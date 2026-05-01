#!/usr/bin/env python3
"""
Generate a PPTX presentation from a JSON slide plan using the EPAM template.
Usage: python generate_pptx.py <slides.json> <output.pptx>

Supported layouts (use these exact strings in JSON):
  - Cover_Dark                 : title slide (title, subtitle, label)
  - Contents                   : table of contents (items list of {number, text} pairs, up to 6)
  - Subsection_Divider         : section break (title, body, subtitle)
  - Quote_Large_Black          : big quote/impact slide (title, quote, label)
  - Title+Small_Text           : content slide with bullets + optional summary line (title, body list, subtitle)
  - Title Only                 : title + short section label only (title, section_label) — no body area
  - Our_People_1_Person_White  : story/example slide (section_title, name, job_title, body)
  - Thank_you                  : closing slide (name, job_title, email)
"""

import json
import sys
from pathlib import Path
from pptx import Presentation

SCRIPT_DIR = Path(__file__).parent
TEMPLATE_PATH = SCRIPT_DIR.parent / "assets" / "template.pptx"

LAYOUT_INDEX = {
    "Cover_Dark": 1,
    "Contents": 3,
    "Subsection_Divider": 12,
    "Quote_Large_Black": 26,
    "Title Only": 9,
    "Title+Small_Text": 31,
    "Our_People_1_Person_White": 21,
    "Thank_you": 68,
}


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_text(ph, text):
    if ph is None or not text:
        return
    tf = ph.text_frame
    tf.clear()
    tf.paragraphs[0].text = str(text)


def set_bullets(ph, items):
    if ph is None or not items:
        return
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0


def remove_all_slides(prs):
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        rId = sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)


def add_slide(prs, slide_def):
    layout_name = slide_def.get("layout", "Title+Subtitle")
    layout_idx = LAYOUT_INDEX.get(layout_name, 10)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    title   = slide_def.get("title", "")
    body    = slide_def.get("body", [])
    subtitle = slide_def.get("subtitle", "")
    quote   = slide_def.get("quote", "")
    label   = slide_def.get("label", "")

    if layout_name == "Cover_Dark":
        # idx=10: main title (large, top), idx=11: subtitle/description, idx=12: date label (bottom)
        set_text(get_ph(slide, 10), title)
        set_text(get_ph(slide, 11), subtitle)
        set_text(get_ph(slide, 12), label)

    elif layout_name == "Contents":
        # idx=0: "Contents" title
        # idx=10–15: row number labels (6 rows)
        # idx=16–21: row section-name labels (6 rows, paired with numbers)
        # items: list of {number, text} dicts, up to 6
        set_text(get_ph(slide, 0), title or "Contents")
        items = slide_def.get("items", [])
        for i, item in enumerate(items[:6]):
            set_text(get_ph(slide, 10 + i), str(item.get("number", i + 1)))
            set_text(get_ph(slide, 16 + i), item.get("text", ""))

    elif layout_name == "Subsection_Divider":
        set_text(get_ph(slide, 0), title)
        if isinstance(body, list):
            set_bullets(get_ph(slide, 10), body)
        else:
            set_text(get_ph(slide, 10), body)
        set_text(get_ph(slide, 16), subtitle)

    elif layout_name == "Quote_Large_Black":
        # idx=11: small section title (top-left), idx=0: large quote body (center), idx=16: attribution (bottom)
        set_text(get_ph(slide, 11), title)
        set_text(get_ph(slide, 0), quote)
        set_text(get_ph(slide, 16), label)

    elif layout_name == "Title Only":
        # idx=0: full-width title bar, idx=28: small section label (0.25" tall, top-left — short text only)
        # NOTE: this layout has NO large body area; use Title+Small_Text for bullet/paragraph content
        set_text(get_ph(slide, 0), title)
        set_text(get_ph(slide, 28), slide_def.get("section_label", ""))

    elif layout_name == "Title+Small_Text":
        # idx=28: small section label (top), idx=0: title, idx=26: bold summary line, idx=27: body/bullets
        set_text(get_ph(slide, 0), title)
        if isinstance(body, list):
            set_bullets(get_ph(slide, 27), body)
        else:
            set_text(get_ph(slide, 27), body)
        set_text(get_ph(slide, 26), subtitle)

    elif layout_name == "Our_People_1_Person_White":
        # idx=11: section title (~3.8" wide, single line, top-left)
        # idx=25: person name (~3.8" wide x 1.3" tall — use for name + optional subtitle line)
        # idx=26: role/title (~5.1" wide x 0.24" tall — one short line only)
        # idx=16: narrative body (~5.1" wide x 2.8" tall — 3–5 sentences max)
        # idx=27: picture placeholder (right half of slide — cannot be filled via python-pptx)
        set_text(get_ph(slide, 11), slide_def.get("section_title", ""))
        set_text(get_ph(slide, 25), slide_def.get("name", ""))
        set_text(get_ph(slide, 26), slide_def.get("job_title", ""))
        if isinstance(body, list):
            set_bullets(get_ph(slide, 16), body)
        else:
            set_text(get_ph(slide, 16), body)

    elif layout_name == "Thank_you":
        # idx=0: "Thank you!" title, idx=12: contact label, idx=13: name, idx=14: job title, idx=15: email
        set_text(get_ph(slide, 0), "Thank you!")
        set_text(get_ph(slide, 13), slide_def.get("name", ""))
        set_text(get_ph(slide, 14), slide_def.get("job_title", ""))
        set_text(get_ph(slide, 15), slide_def.get("email", ""))


def main():
    if len(sys.argv) < 3:
        print("Usage: generate_pptx.py <slides.json> <output.pptx>")
        sys.exit(1)

    slides_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    with open(slides_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    for slide_def in data.get("slides", []):
        add_slide(prs, slide_def)

    prs.save(str(output_path))
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
